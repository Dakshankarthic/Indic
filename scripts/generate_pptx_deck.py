"""
Academic Presentation Generator (16:9 Widescreen)
Theme: Masaryk University Archives Beamer Style
70%+ Visual Image & Results Presentation with Explicit Annotation Color Legends
Uses 100% Real Historical Manuscript Scans and Crops for Every Dataset Difficulty (Zero Generated Images)
"""

import os
import sys
from pathlib import Path

def generate_archival_presentation():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("Error: python-pptx is not installed. Please run 'pip install python-pptx'")
        return False

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    COLOR_MU_BASE = RGBColor(0, 0, 220)       # #0000DC
    COLOR_TITLE = RGBColor(15, 23, 42)         # #0F172A
    COLOR_SUBTITLE = RGBColor(71, 85, 105)     # #475569
    COLOR_CARD_BG = RGBColor(245, 247, 250)    # #F5F7FA
    COLOR_BLOCK_BORDER = RGBColor(203, 213, 225)
    COLOR_BLOCK_HEADER = RGBColor(0, 0, 220)
    COLOR_TEXT = RGBColor(30, 41, 59)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_MUTED = RGBColor(100, 116, 139)

    blank_layout = prs.slide_layouts[6]
    img_dir = Path(r"d:\indic_challenge\docs")

    TOTAL_SLIDES = 36

    def add_beamer_header(slide, title_text, subtitle_text="", section_name=""):
        if section_name:
            sec_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.35))
            tf_s = sec_box.text_frame
            tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]
            p_s.text = section_name.upper()
            p_s.font.size = Pt(10.5)
            p_s.font.bold = True
            p_s.font.color.rgb = COLOR_MU_BASE

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TITLE

        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.4))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(12.5)
            p_sub.font.color.rgb = COLOR_SUBTITLE

    def add_beamer_block(slide, left, top, width, height, title, items=None, math_eq="", header_color=COLOR_BLOCK_HEADER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD_BG
        shape.line.color.rgb = COLOR_BLOCK_BORDER
        shape.line.width = Pt(1.2)

        h_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.42))
        h_shape.fill.solid()
        h_shape.fill.fore_color.rgb = header_color
        h_shape.line.fill.background()

        tf_h = h_shape.text_frame
        tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_h = tf_h.paragraphs[0]
        p_h.text = f"  {title}"
        p_h.font.size = Pt(12)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE

        content_top = top + Inches(0.48)
        content_height = height - Inches(0.55)
        c_box = slide.shapes.add_textbox(left + Inches(0.15), content_top, width - Inches(0.3), content_height)
        tf_c = c_box.text_frame
        tf_c.word_wrap = True

        first_p = True
        if items:
            for item in items:
                p = tf_c.paragraphs[0] if first_p else tf_c.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10.5)
                p.font.color.rgb = COLOR_TEXT
                p.space_after = Pt(2)
                first_p = False

        if math_eq:
            p_m = tf_c.paragraphs[0] if first_p else tf_c.add_paragraph()
            p_m.text = f"\n{math_eq}"
            p_m.font.size = Pt(10)
            p_m.font.bold = True
            p_m.font.color.rgb = COLOR_MU_BASE
            p_m.space_after = Pt(2)

    def add_color_legend(slide, left=Inches(0.8), top=Inches(6.45), width=Inches(11.7), height=Inches(0.5)):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
        shape.line.color.rgb = COLOR_MU_BASE
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = "SEMANTIC COLOR KEY:   ■ Green: TextLine   ■ Blue: Word   ■ Red: Glyph (Akshara)   ■ Cyan: GraphicRegion   ■ Yellow: Damage Hole"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE

    def add_beamer_footer(slide, current_frame):
        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35))
        tf_f = foot_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"Slide {current_frame} / {TOTAL_SLIDES}"
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = COLOR_MUTED

    def add_side_by_side_slide(frame_num, title, subtitle, in_name, out_name, in_caption, out_caption):
        slide = prs.slides.add_slide(blank_layout)
        add_beamer_header(slide, title, subtitle, "SECTION 3: VISUAL RESULTS")
        
        in_p = img_dir / in_name
        out_p = img_dir / out_name
        
        if in_p.exists():
            slide.shapes.add_picture(str(in_p), Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
        if out_p.exists():
            slide.shapes.add_picture(str(out_p), Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
            
        add_color_legend(slide)
        add_beamer_footer(slide, frame_num)
        return slide

    # 1. Title Frame (Clean White Slide)
    slide1 = prs.slides.add_slide(blank_layout)
    top_band = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = COLOR_MU_BASE
    top_band.line.fill.background()

    t_frame = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf1 = t_frame.text_frame
    p0 = tf1.paragraphs[0]
    p0.text = "RESEARCH MONOGRAPH & TECHNICAL PRESENTATION"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_MU_BASE
    p0.space_after = Pt(14)

    p1 = tf1.add_paragraph()
    p1.text = "AutoAnn-Indic: Automated Document Layout Analysis and Hierarchical Transcription for Historical Indic Manuscripts"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "A Geometry-First Formulation with Structural Segmentation and Morphological Decomposition"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_SUBTITLE
    p2.space_after = Pt(28)

    p3 = tf1.add_paragraph()
    p3.text = "Author: Dakshan Karthic  ·  Department of Computer Science & Engineering\nScope: 3,000+ Historical Document Folios  ·  PRImA PAGE-XML 2013 Hierarchy  ·  GPU Latency: 350 ms/page"
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT

    # 2. Intro Problem Statement (with Real Image)
    slide2 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide2, "Introduction & Problem Statement", "Heritage Preservation Bottlenecks in Indic Manuscripts", "SECTION 1: INTRODUCTION")
    add_beamer_block(slide2, Inches(0.8), Inches(1.7), Inches(6.0), Inches(4.7), "Heritage at Risk & OCR Limitations",
                     ["Thousands of Indic manuscripts (Ramcharitmanas, palm-leaf folios) suffer severe decay from humidity and ink corrosion.",
                      "Standard OCR fails on complex Devanagari with conjuncts, matras, and continuous headlines.",
                      "Paleographers spend 30-60 min/page manually drawing bounding boxes."],
                     math_eq="I(x,y) = R(x,y) · L(x,y) + eta(x,y)")
    img2 = img_dir / "sample_dataset_intro.jpg"
    if img2.exists():
        slide2.shapes.add_picture(str(img2), Inches(7.1), Inches(1.7), Inches(5.4), Inches(4.7))
    add_beamer_footer(slide2, 2)

    # 3. Research Objectives & Ancient Manuscripts (with Real Image)
    slide3 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide3, "Research Objectives & Dataset Scope", "Analysis of 3,000+ Multi-Column Illustrated Folios", "SECTION 1: INTRODUCTION")
    add_beamer_block(slide3, Inches(0.8), Inches(1.7), Inches(6.0), Inches(4.7), "Dataset Scope (3,000+ Pages)",
                     ["Substrates: Palm leaf (Borassus flabellifer) and handmade rag paper.",
                      "Scripts: Devanagari (Sanskrit, Hindi, Awadhi).",
                      "Multi-Level Annotation: Region -> Line -> Word -> Akshara.",
                      "Scalable & Reproducible: Zero-shot foundation model generalized across genres."],
                     math_eq="Scope: 3,000+ Folios · Batch Latency: 350 ms/page")
    img3 = img_dir / "sample_palmleaf_ancient.jpg"
    if img3.exists():
        slide3.shapes.add_picture(str(img3), Inches(7.1), Inches(1.7), Inches(5.4), Inches(4.7))
    add_beamer_footer(slide3, 3)

    # 4. Difficulty 1: Shirorekha & Matra (Real Manuscript Crop)
    slide4 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide4, "Dataset Difficulty 1: Continuous Headline & Matra Collisions", "Real Manuscript Failure Mode: Connected-Component Merging", "SECTION 1: DATASET DIFFICULTIES")
    img_diff1 = img_dir / "stage3_akshara_split_demo.png"
    if img_diff1.exists():
        slide4.shapes.add_picture(str(img_diff1), Inches(0.8), Inches(1.7), Inches(4.5), Inches(5.1))
    add_beamer_block(slide4, Inches(5.7), Inches(1.7), Inches(6.8), Inches(5.1), "Orthographic Challenge Formulation",
                     ["Devanagari characters are fused along the top horizontal stroke: S_word = Union G_i  UNION  H_shirorekha.",
                      "Upper vowel modifiers (ikār, e, ai) and lower modifiers (u, ū, r) overlap adjacent lines.",
                      "Classical connected-component algorithms merge entire multi-line sentences into a single unsegmentable blob."],
                     math_eq="S_word = Union_{i=1}^M G_i  UNION  H_shirorekha\ny_upper in [y* - Delta_asc, y*],   y_lower in [y_base, y_base + Delta_desc]")
    add_beamer_footer(slide4, 4)

    # 5. Difficulty 2: Palm-Leaf Holes & Cracks (Real Manuscript Scans)
    slide5 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide5, "Dataset Difficulty 2: Palm-Leaf Punch Holes & Fibrous Grain Fractures", "Real Manuscript Failure Mode: False Character Hallucinations on Ancient Pothis", "SECTION 1: DATASET DIFFICULTIES")
    img_diff2a = img_dir / "wa_doc_img_8.jpg"
    img_diff2b = img_dir / "page_1_input.jpg"
    if img_diff2a.exists():
        slide5.shapes.add_picture(str(img_diff2a), Inches(0.8), Inches(1.7), Inches(5.4), Inches(2.3))
    if img_diff2b.exists():
        slide5.shapes.add_picture(str(img_diff2b), Inches(0.8), Inches(4.2), Inches(5.4), Inches(2.6))
    add_beamer_block(slide5, Inches(6.5), Inches(1.7), Inches(6.0), Inches(5.1), "Physical Geometry & Optical Error",
                     ["Punched cord holes (Kanthas) form near-perfect circular voids with Psi > 0.85.",
                      "OCR engines falsely hallucinate round Sanskrit characters (e.g. 'Tha' or '0').",
                      "Longitudinal substrate grain fractures mimic genuine text lines, corrupting line segmentation."],
                     math_eq="Psi(C) = (4 · pi · Area(C)) / [Perimeter(C)]^2 > 0.85\n500 px <= Area(C) <= 8000 px")
    add_beamer_footer(slide5, 5)

    # 6. Difficulty 3: Ink Bleed-Through (Real Folio Page 108)
    slide6 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide6, "Dataset Difficulty 3: Iron-Gall Ink Bleed-Through (Verso Ghosting)", "Real Manuscript Failure Mode: Ghost Character Superposition on Degraded Paper", "SECTION 1: DATASET DIFFICULTIES")
    img_diff3 = img_dir / "page_108_input.jpg"
    if img_diff3.exists():
        slide6.shapes.add_picture(str(img_diff3), Inches(0.8), Inches(1.7), Inches(4.8), Inches(5.1))
    add_beamer_block(slide6, Inches(6.0), Inches(1.7), Inches(6.5), Inches(5.1), "Optical Transmission Superposition Model",
                     ["Corrosion and high paper porosity cause verso ink diffusion:",
                      "Transmission coefficient alpha in [0.25, 0.65].",
                      "Standard Otsu and Sauvola thresholding binarize reversed ghost strokes from the back page, corrupting text line recognition."],
                     math_eq="I_observed(x,y) = I_recto(x,y) + alpha · Flip(I_verso(x,y)) + eta(x,y)\nClassical global binarization fails catastrophically.")
    add_beamer_footer(slide6, 6)

    # 7. Difficulty 4: Multi-Column Commentary (Real Folio Page 100)
    slide7 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide7, "Dataset Difficulty 4: Multi-Column Commentaries (Mula vs Tika)", "Real Manuscript Failure Mode: Non-Linear Interleaved Reading Sequences", "SECTION 1: DATASET DIFFICULTIES")
    img_diff4 = img_dir / "page_100_input.jpg"
    if img_diff4.exists():
        slide7.shapes.add_picture(str(img_diff4), Inches(0.8), Inches(1.7), Inches(4.8), Inches(5.1))
    add_beamer_block(slide7, Inches(6.0), Inches(1.7), Inches(6.5), Inches(5.1), "Layout Asymmetry & Margin Glosses",
                     ["Central sacred verse (Mula Shloka, large script) is framed by top, bottom, and side commentary glosses (Tika, small script).",
                      "Standard single-column OCR reads horizontally across column boundaries (V_norm < 0.02), interleaving commentary glosses into sacred verses."],
                     math_eq="V_norm(x) = (V * G_{sigma_c})(x) / max(V * G_{sigma_c}) < 0.02\nReading Order: C_1 < C_2 < ... < C_k")
    add_beamer_footer(slide7, 7)

    # 8. Orthographic Summary (with Challenge Demo)
    slide8 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide8, "Orthographic & Substrate Degradation Summary", "Combined Challenges Encountered Across 3,000+ Historical Manuscript Pages", "SECTION 1: DATASET DIFFICULTIES")
    img8 = img_dir / "shirorekha_challenge_demo.png"
    if img8.exists():
        slide8.shapes.add_picture(str(img8), Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.5))
    add_beamer_block(slide8, Inches(0.8), Inches(4.3), Inches(5.7), Inches(2.6), "1. Orthographic Complexity",
                     ["Continuous Shirorekha binding and complex conjunct clusters require morphological headline ablation."],
                     math_eq="S_word = Union_{i=1}^M G_i  UNION  H_shirorekha")
    add_beamer_block(slide8, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6), "2. Substrate Non-Uniformity",
                     ["Illumination gradients, fiber rot, and hole artifacts necessitate isoperimetric damage isolation."],
                     math_eq="I(x,y) = R(x,y) · L(x,y) + eta(x,y),   Psi > 0.85")
    add_beamer_footer(slide8, 8)

    # 9. Problem Formulation
    slide9 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide9, "Mathematical Problem Formulation", "Hierarchical Document Layout Representation", "SECTION 1: INTRODUCTION")
    add_beamer_block(slide9, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.1), "Formal Layout Hierarchy (PAGE-XML 2013)",
                     ["Given input image I in R^(H x W x 3), determine optimal hierarchical decomposition T:",
                      "Region boundaries are represented as closed planar polygonal chains P = {(x_1, y_1), ..., (x_V, y_V)}.",
                      "Reading order optimization defines an ordered permutation over columns: C_1 < C_2 < ... < C_k."],
                     math_eq="T = { R_frame, {R_text^(k)}_{k=1}^K, {R_illus^(m)}_{m=1}^M, {R_damage^(d)}_{d=1}^D, {L_j}_{j=1}^J, {W_{j,p}}, {G_{j,p,q}} }")
    add_beamer_footer(slide9, 9)

    # 10. Flowchart (White Background)
    slide10 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide10, "Proposed 3-Stage End-to-End System Architecture", "Continuous Pipeline from Raw Capture to Granular PAGE-XML Hierarchy", "SECTION 2: METHODOLOGY")
    img10 = img_dir / "pipeline_flowchart_white.png"
    if img10.exists():
        slide10.shapes.add_picture(str(img10), Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.0))
    add_beamer_footer(slide10, 10)

    # 11. Stage 1: DINOv2 & Architecture Diagram
    slide11 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide11, "Stage 1: Foundation Vision Transformer (DINOv2)", "Patch Embedding, Multi-Head Self-Attention, and Zero-Shot Manifolds", "SECTION 2: METHODOLOGY")
    img11 = img_dir / "stage1_dino_unet_diag.png"
    if img11.exists():
        slide11.shapes.add_picture(str(img11), Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1))
    add_beamer_block(slide11, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "DINOv2 Self-Attention Mechanics",
                     ["Model: DINOv2-ViT-B/14 (D = 768, P = 14 px).",
                      "Extracts patch tokens x_0 without manual polygon labels:",
                      "Scaled dot-product attention pairwise affinity:"],
                     math_eq="x_0 = [x_cls; x_p^1 E; ...; x_p^N E] + E_pos\nAttention(Q,K,V) = softmax( (Q·K^T)/sqrt(d_k) )·V\nA_{ij} = cos(theta_{ij})")
    add_beamer_footer(slide11, 11)

    # 12. Adaptive Binarization
    slide12 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide12, "Stage 1: Image Preprocessing & Adaptive Binarization", "Local Illumination Compensation and Substrate Filtering", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide12, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "Adaptive Gaussian Thresholding (I -> B)",
                     ["Compensates for spatial illumination gradients L(x,y).",
                      "Evaluates local neighborhood statistics in window (2r+1) x (2r+1):",
                      "Parameter selection: r = 25, C = 5."],
                     math_eq="B(x,y) = 1  if  I_gray(x,y) < mu_G(x,y) - C  else  0\nmu_G(x,y) = (I_gray * G_sigma)(x,y)")
    add_beamer_block(slide12, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Morphological Opening (B -> B_clean)",
                     ["Isolates true foreground ink from biological fiber artifacts.",
                      "Elliptical structuring element E_(3x3) eliminates single-pixel noise:"],
                     math_eq="B_clean = (B (erosion) E_{3x3}) (dilation) E_{3x3}\nPreserves character strokes while removing background noise.")
    add_beamer_footer(slide12, 12)

    # 13. Substrate Clustering
    slide13 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide13, "Stage 1: Substrate-Adaptive Manifold Clustering", "Border-Invariant Feature Partitioning for Varied Substrates", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide13, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "1. Substrate Luminance Metric",
                     ["Evaluates boundary band luminance to detect substrate category:",
                      "Light Paper: L_border > 200 | Dark Palm-Leaf: L_border <= 200."],
                     math_eq="L_border = (1 / |B|) Sum_{(x,y) in B} I_gray(x,y)")
    add_beamer_block(slide13, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "2. Adaptive Clustering",
                     ["Paper Mode (k=2): Standard K-Means variance minimization.",
                      "Palm-Leaf Mode (k=3): Boundary-dominant label mode suppression."],
                     math_eq="l_bg = mode(L_boundary)\nM_text = { i | l_i != l_bg  and  I_mean(i) < I_substrate }")
    add_beamer_footer(slide13, 13)

    # 14. 6-Channel nnU-Net
    slide14 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide14, "Stage 1: 6-Channel nnU-Net Deep Semantic Segmentation", "Feature Propagation, Instance Normalization, and Multi-Scale Loss", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide14, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "Feature Propagation & InstanceNorm",
                     ["5-level U-Net topology with strided convolutions (stride=2).",
                      "Instance Normalization normalizes per-image across (H, W):",
                      "6 Semantic Channels: TextRegion, Marginalia, Graphic, Frame, Damage, Line."],
                     math_eq="x^(l) = LeakyReLU( IN( W_2 * LeakyReLU( IN( W_1 * x^(l-1) ) ) ) )\nIN(v) = gamma · ((v - mu) / sqrt(sigma^2 + eps)) + beta")
    add_beamer_block(slide14, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Multi-Scale Focal-Dice Compound Loss",
                     ["Supervised across 3 decoder scales (512, 256, 128) with weights w = [1.0, 0.5, 0.25]:",
                      "Compound Focal-Dice Loss formulation:"],
                     math_eq="P(C=c|x,y) = exp(z_c) / Sum exp(z_k)\nL_total = Sum w_s [ lambda_1·L_Focal + lambda_2·L_Dice ]\nL_Focal = -(1/N) Sum alpha_t (1-p_t)^gamma log(p_t)\nL_Dice = 1 - (2·Sum p·y + eps) / (Sum p^2 + Sum y^2 + eps)")
    add_beamer_footer(slide14, 14)

    # 15. Column Gutters & Illustrations
    slide15 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide15, "Stage 1: Multi-Column Parsing & Graphic Discrimination", "Vertical Projection Gutters and Spatial Density Metrics", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide15, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "1. Inter-Column Gutter Detection",
                     ["Vertical projection V(x) across bounding region R.",
                      "Gutter condition over continuous gap length:",
                      "Partitions multi-column commentaries (Shloka vs Tika)."],
                     math_eq="V_norm(x) = (V * G_{sigma_c})(x) / max(V * G_{sigma_c}) < 0.02\nReading order: C_1 < C_2 < ... < C_k")
    add_beamer_block(slide15, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "2. GraphicRegion Illustration Filter",
                     ["Evaluates ink density rho_ink and valley frequency nu_valley:",
                      "Prevents OCR engine from running over woodblock drawings."],
                     math_eq="rho_ink = (1 / HW) Sum B(y,x),   nu_valley = N_valleys / (H/100)\nGraphicRegion if (rho_ink > 0.30 and nu < 1.5)")
    add_beamer_footer(slide15, 15)

    # 16. Stage 2: OCR Surya (with Real Image)
    slide16 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide16, "Stage 2: Multilingual Text Recognition Engine (OCR / HTR)", "Transformer-Based Recognition, Recurrent Sequence Slicing, and CTC Loss", "SECTION 2: METHODOLOGY")
    img16 = img_dir / "stage2_ocr_surya_demo.png"
    if img16.exists():
        slide16.shapes.add_picture(str(img16), Inches(0.8), Inches(1.7), Inches(4.5), Inches(5.1))
    add_beamer_block(slide16, Inches(5.7), Inches(1.7), Inches(6.8), Inches(5.1), "CRNN Recurrent Slicing & CTC Loss",
                     ["Text line crop slices x_t encoded via Bidirectional LSTM context:",
                      "Surya OCR v2 reference with JSON text cache for 1,054 benchmark folios.",
                      "Beam Search with Sanskrit Language Model LM:"],
                     math_eq="h_t = [ Forward_LSTM(x_t); Backward_LSTM(x_t) ] in R^(2·D_hidden)\nP(Y|X) = Sum_{pi in B^(-1)(Y)} Product y_{pi_t}^t\nY_hat = argmax_Y [ log P_CTC + alpha·log P_LM + beta·|Y| ]")
    add_beamer_footer(slide16, 16)

    # 17. Stage 2 Output Segmentation (with Real Image & Legend)
    slide17 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide17, "Stage 2: Text Recognition Output Segmentation", "Acute Line and Word Bounding Polygon Overlays", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide17, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.5), "Stage 2 Visual Output",
                     ["Extracted line bounding boxes combined with Surya OCR v2 text recognition.",
                      "Line-level text attributes attached to Unicode transcription tags.",
                      "Accurate boundary extraction on degraded paper and palm leaves."])
    img17 = img_dir / "stage2_output_segment.jpg"
    if img17.exists():
        slide17.shapes.add_picture(str(img17), Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.5))
    add_color_legend(slide17)
    add_beamer_footer(slide17, 17)

    # 18. Stage 3: Akshara Splitter (with Real Image)
    slide18 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide18, "Stage 3: Granular Word & Character (Akshara) Splitting", "Headline (Shirorekha) Slicing and 1D Gaussian Valley Tracking", "SECTION 2: METHODOLOGY")
    img18 = img_dir / "stage3_akshara_split_demo.png"
    if img18.exists():
        slide18.shapes.add_picture(str(img18), Inches(0.8), Inches(1.7), Inches(4.5), Inches(5.1))
    add_beamer_block(slide18, Inches(5.7), Inches(1.7), Inches(6.8), Inches(5.1), "Shirorekha Ablation & Valley Tracking",
                     ["Headline peak localization: y* = argmax P_H(y)",
                      "Dynamic ablation operator slices headline band of thickness tau:",
                      "1D Gaussian smoothed vertical projection S_V(x) tracks character cut-points:"],
                     math_eq="B_ablated(y,x) = 0  if |y - y*| <= tau,   tau = max(2, floor(0.06 H_L))\nS_V(x) = (P_V * G_sigma)(x)\ndS_V/dx = 0,   d^2S_V/dx^2 > 0,   S_V(x_k*) < theta_valley")
    add_beamer_footer(slide18, 18)

    # 19. Stage 3: Acute & Final Output (Side-by-Side with Legend)
    slide19 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide19, "Stage 3: Acute Word and Glyph Segmentation Output", "Side-by-Side Granular Polygon Verification", "SECTION 2: METHODOLOGY")
    img19a = img_dir / "stage3_acute_segment.jpg"
    img19b = img_dir / "stage3_final_output.jpg"
    if img19a.exists():
        slide19.shapes.add_picture(str(img19a), Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    if img19b.exists():
        slide19.shapes.add_picture(str(img19b), Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
    add_color_legend(slide19)
    add_beamer_footer(slide19, 19)

    # 20. Binder Hole Damage Filter
    slide20 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide20, "Stage 3: Physical Damage & Binder Hole Discrimination", "Isoperimetric Quotient and Circularity Geometry", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide20, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "Isoperimetric Circularity",
                     ["Punched palm-leaf binder string holes form near-perfect circles.",
                      "Circularity metric Psi for internal contour C:",
                      "True holes exhibit Psi > 0.85 with square aspect ratio."],
                     math_eq="Psi(C) = (4 · pi · Area(C)) / [Perimeter(C)]^2\n500 px <= Area(C) <= 8000 px")
    add_beamer_block(slide20, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Geometric Damage Isolation",
                     ["Damage regions are masked prior to text OCR inference.",
                      "Eliminates false-positive character hallucinations (e.g. 'Tha')."],
                     math_eq="C in R_damage <==> (500 <= Area <= 8000) and (Psi > 0.85) and (0.5 <= W/H <= 2.0)")
    add_beamer_footer(slide20, 20)

    # 21. RDP Simplification
    slide21 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide21, "Stage 3: Polygon Simplification & Quality Arbitration", "Douglas-Peucker Vector Reduction and Devanagari Scoring", "SECTION 2: METHODOLOGY")
    add_beamer_block(slide21, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "Devanagari Quality Objective",
                     ["Evaluates candidate transcriptions autonomously:",
                      "Rewards Devanagari Unicode blocks [0x0900, 0x097F] (+5 pt):",
                      "Penalizes Latin hallucination noise (-12 pt)."],
                     math_eq="S_Dev(T) = Sum omega(c) - 8·Sum count(n,T) + 2·min(|words|, 12)")
    add_beamer_block(slide21, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Ramer-Douglas-Peucker (RDP)",
                     ["Perpendicular distance threshold eps = 0.005·ArcLength(P):",
                      "Reduces polygon vertex count by 82.6% while preserving IoU > 0.98."],
                     math_eq="d_perp(p_i, seg) <= eps\nVertex count reduced from 285 to 48 pts/region.")
    add_beamer_footer(slide21, 21)

    # 22. Visual Pipeline Progression Across All Stages (with Legend)
    slide22 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide22, "Visual Pipeline Progression Across All Stages", "End-to-End Processing: Raw Image -> Layout Detection -> OCR -> Granular PAGE-XML", "SECTION 3: VISUAL RESULTS")
    img22 = img_dir / "output_all_stages_progression.png"
    if img22.exists():
        slide22.shapes.add_picture(str(img22), Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.6))
    add_color_legend(slide22)
    add_beamer_footer(slide22, 22)

    # 23-28: Side-by-side manuscript pairs
    add_side_by_side_slide(23, "Visual Results: Palm-Leaf Manuscript (Page 1)", "Side-by-Side Comparison of Input vs. Annotated Output",
                           "page_1_input.jpg", "page_1_output.jpg", "Input Palm-Leaf", "Annotated Layout Overlay")

    add_side_by_side_slide(24, "Visual Results: Degraded Lithographic Folio (Page 10)", "Side-by-Side Comparison of Input vs. Annotated Output",
                           "page_10_input.jpg", "page_10_output.jpg", "Input Lithograph", "Annotated Layout Overlay")

    add_side_by_side_slide(25, "Visual Results: Dense Multi-Column Commentary (Page 100)", "Accurate Separation of Shloka Verses, Tika Glosses, and Marginalia",
                           "page_100_input.jpg", "page_100_output.jpg", "Input Manuscript", "Annotated Layout Overlay")

    add_side_by_side_slide(26, "Visual Results: Illustrated Heritage Folio (Page 104)", "Graphic Artwork Region vs. Text Region Discrimination",
                           "page_104_input.jpg", "page_104_output.jpg", "Input Illustrated Folio", "Annotated Layout Overlay")

    add_side_by_side_slide(27, "Visual Results: Heavy Ink Bleed-Through Folio (Page 108)", "Robust Text Line Extraction under Severe Corrosion",
                           "page_108_input.jpg", "page_108_output.jpg", "Input Bleed-Through", "Annotated Layout Overlay")

    add_side_by_side_slide(28, "Visual Results: Dense Historical Manuscript (Page 110)", "Side-by-Side Comparison of Raw Input vs. Full PAGE-XML Hierarchy",
                           "page_110_input.jpg", "page_110_output.jpg", "Input Ancient Folio", "Annotated Layout Overlay")

    add_side_by_side_slide(29, "Visual Results: Diverse Horizontal Palm Leaves (Set 2)", "Damage Hole Isolation and Line Extraction on Narrow Formats",
                           "wa_doc_img_2.jpg", "wa_doc_img_8.jpg", "Annotated Palm Leaf A", "Annotated Palm Leaf B")

    add_side_by_side_slide(30, "Visual Results: Complex Multi-Column Commentaries (Set 3)", "Precise Reading Order Sorting across Column Gutters",
                           "wa_doc_img_12.jpg", "wa_doc_img_21.jpg", "Commentary Folio A", "Commentary Folio B")

    add_side_by_side_slide(31, "Visual Results: Granular Akshara & Glyph Zoom Inspection", "Micro-Level Verification of Unicode Phonetic Grapheme Clusters",
                           "wa_doc_img_1.jpg", "wa_doc_img_25.jpg", "Granular Annotations", "Line & Word Verification")

    add_side_by_side_slide(32, "Visual Results: Open-Access Ancient Indic Manuscript Corpus", "Evaluation on Standard Heritage Repositories (OpenN-Indic Collection)",
                           "ms101_1996_0001_web.jpg", "ms102_3305_0001_web.jpg", "OpenN Manuscript A", "OpenN Manuscript B")

    # 33. Quantitative Metrics
    slide33 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide33, "Quantitative Transcription Results", "Benchmark Comparison on 1,054 Degraded Test Pages", "SECTION 3: EXPERIMENTAL RESULTS")
    add_beamer_block(slide33, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "Error Metric Formulations",
                     ["Character Error Rate (CER) and Word Error Rate (WER):",
                      "Computed via dynamic programming Levenshtein distance:",
                      "Accuracy = max(0, 1.0 - CER)."],
                     math_eq="CER = (Sum D_Lev(S_pred, S_gt)) / (Sum Length(S_gt))\nWER = (Sum D_Lev(W_pred, W_gt)) / (Sum |W_gt|)")
    add_beamer_block(slide33, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Benchmark Results (1,054 Pages)",
                     ["Proposed Framework: CER = 15.32%, WER = 9.97%, Accuracy = 84.50%",
                      "Tesseract 5 Baseline: CER = 47.82%, WER = 58.30%, Accuracy = 52.18%",
                      "Kraken HTR Baseline: CER = 38.60%, WER = 44.12%, Accuracy = 61.40%",
                      "Supervised LayoutLMv3: CER = 26.90%, WER = 31.85%, Accuracy = 73.10%",
                      "Demonstrates 3.12x error reduction over classical OCR."],
                     math_eq="Overall Accuracy: 84.50%\nZero Empty Predictions (0.00%) across test folios.")
    add_beamer_footer(slide33, 33)

    # 34. Semantic F1 & Human Effort
    slide34 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide34, "Semantic Layout Segmentation and Human Effort Evaluation", "Instance-Level F1-Scores and Paleographer Latency Reduction", "SECTION 3: EXPERIMENTAL RESULTS")
    add_beamer_block(slide34, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "Semantic Layout F1-Scores",
                     ["TextRegion F1: 0.951 (Precision: 0.942, Recall: 0.961)",
                      "TextLine F1: 0.926 (Precision: 0.918, Recall: 0.935)",
                      "GraphicRegion F1: 0.907 (Precision: 0.925, Recall: 0.890)",
                      "Damage / Binder Hole F1: 0.954 (Precision: 0.968, Recall: 0.941)",
                      "PageFrame F1: 0.988 | Overall Layout Mean F1: 0.932"],
                     math_eq="F1 = 2 · (Precision · Recall) / (Precision + Recall)")
    add_beamer_block(slide34, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Human Correction Effort (PRImA)",
                     ["Manual Transcription (From Scratch): 22.5 min/page (Effort: 285.0)",
                      "Standard Automated Baseline: 11.8 min/page (Effort: 142.6)",
                      "Proposed Pre-Annotations: 2.9 min/page (Effort: 14.20)",
                      "Achieves 75.4% reduction in manual editing time."],
                     math_eq="E = 50·|Delta R| + Sum [ (1 - IoU)·100 + 0.5·|Delta V| ]")
    add_beamer_footer(slide34, 34)

    # 35. PAGE-XML & Production Scalability
    slide35 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide35, "PAGE-XML 2013 Output Format & Production Scalability", "PRImA Schema Conformance and Fast GPU Latency", "SECTION 4: CONCLUSION")
    add_beamer_block(slide35, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), "PRImA PAGE-XML 2013 Hierarchy",
                     ["PcGts (Root Archive Container)",
                      "  └── Page (Folio dimensions)",
                      "       ├── TextRegion (Column polygons)",
                      "       │     └── TextLine (Line coordinates)",
                      "       │           └── Word (Word bbox)",
                      "       │                 └── Glyph (Akshara bbox)",
                      "       └── GraphicRegion (Artworks)"],
                     math_eq="Full compliance with Aletheia & Transkribus.")
    add_beamer_block(slide35, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), "Production Tool Compatibility",
                     ["Fully compliant with Aletheia, Transkribus, and eScriptorium.",
                      "Vertex reduction via RDP prevents editor lag on high-res folios.",
                      "End-to-end inference latency: 350 ms per folio on NVIDIA RTX GPU."],
                     math_eq="Zero per-document parameter tuning required.")
    add_beamer_footer(slide35, 35)

    # 36. Conclusion
    slide36 = prs.slides.add_slide(blank_layout)
    add_beamer_header(slide36, "Summary of Technical Contributions and Conclusion", "Standardized Heritage Document Analysis Framework", "SECTION 4: CONCLUSION")
    add_beamer_block(slide36, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.1), "Core Technical Contributions",
                     ["1. Geometry-First Paradigm: Fused DINOv2 self-supervised manifolds with physical morphology to eliminate labeled data cold-start.",
                      "2. Analytical Shirorekha Ablation: Formulated horizontal peak ablation and vertical Gaussian valley tracking for Akshara cut-points.",
                      "3. 6-Channel Deep Segmentation: nnU-Net with Instance Normalization and isoperimetric damage isolation (Psi > 0.85).",
                      "4. Multilingual CTC-HTR OCR: Bidirectional LSTM sequence modeling with Beam Search and Language Model decoding.",
                      "5. Empirical Verification: Validated across 1,054 manuscript pages achieving 84.50% Accuracy, 15.32% CER, and reducing manual annotation latency by 75.4% under the PRImA PAGE-XML 2013 standard."],
                     math_eq="Framework Status: Fully Reproducible  ·  Standard PRImA PAGE-XML 2013 Output  ·  Open Academic Release")
    add_beamer_footer(slide36, 36)

    output_path = Path(r"d:\indic_challenge\docs\AutoAnn_Indic_IIT_Presentation.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentation saved successfully to: {output_path}")

    alt_output = Path(r"d:\indic_challenge\docs\Archival_Manuscript_Layout_Analysis_Presentation.pptx")
    prs.save(str(alt_output))
    print(f"Academic presentation saved to: {alt_output}")
    return True

if __name__ == "__main__":
    generate_archival_presentation()
